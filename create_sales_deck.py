#!/usr/bin/env python3
"""
D365 FO License & Security Optimization Agent - Sales Deck Generator
Generates a professional PowerPoint sales presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette (Azure-inspired dark theme) ──
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)       # Very dark navy
DARK_BG2 = RGBColor(0x1B, 0x25, 0x3B)      # Slightly lighter navy
AZURE_BLUE = RGBColor(0x00, 0x78, 0xD4)     # Microsoft Azure blue
AZURE_LIGHT = RGBColor(0x50, 0xE6, 0xFF)    # Light azure accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xBC, 0xD0)
MID_GRAY = RGBColor(0x6B, 0x7B, 0x93)
GREEN = RGBColor(0x00, 0xC8, 0x53)          # Success green
AMBER = RGBColor(0xFF, 0xB9, 0x00)          # Warning amber
RED_ACCENT = RGBColor(0xFF, 0x4D, 0x4D)     # Alert red
TEAL = RGBColor(0x00, 0xB7, 0xC3)           # Teal accent
PURPLE = RGBColor(0x88, 0x64, 0xD8)         # Purple accent

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_dark_bg(slide, color=DARK_BG):
    """Add a solid dark background to a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE,
                bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    """Add a text box with styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, default_size=16, default_color=LIGHT_GRAY):
    """Add a text box with multiple styled lines. Each line is (text, size, color, bold, alignment)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        color = line_data[2] if len(line_data) > 2 else default_color
        bold = line_data[3] if len(line_data) > 3 else False
        align = line_data[4] if len(line_data) > 4 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Segoe UI"
        p.alignment = align
        p.space_after = Pt(4)
    return txBox


def add_accent_line(slide, left, top, width, color=AZURE_BLUE):
    """Add a horizontal accent line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, body_lines, accent_color=AZURE_BLUE, icon_text=""):
    """Add a styled card with title and body text."""
    card = add_shape(slide, left, top, width, height, DARK_BG2, accent_color, Pt(1.5))
    # Accent bar at top of card
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    # Icon/label
    if icon_text:
        add_textbox(slide, left + Inches(0.2), top + Pt(12), width - Inches(0.4), Inches(0.4),
                    icon_text, font_size=13, color=accent_color, bold=True)
    # Title
    add_textbox(slide, left + Inches(0.2), top + Pt(35), width - Inches(0.4), Inches(0.4),
                title, font_size=16, color=WHITE, bold=True)
    # Body
    y_offset = top + Pt(65)
    for line in body_lines:
        add_textbox(slide, left + Inches(0.2), y_offset, width - Inches(0.4), Inches(0.3),
                    line, font_size=12, color=LIGHT_GRAY)
        y_offset += Pt(20)
    return card


# ════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_dark_bg(slide)

# Large accent shape behind title
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BG
shape.line.fill.background()

# Gradient-like accent bar at top
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Pt(6))
bar.fill.solid()
bar.fill.fore_color.rgb = AZURE_BLUE
bar.line.fill.background()

# Title text
add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
            "D365 FO License & Security", font_size=48, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(2.7), Inches(11), Inches(1),
            "Optimization Agent", font_size=48, color=AZURE_LIGHT, bold=True,
            alignment=PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide, Inches(2), Inches(3.8), Inches(9), Inches(0.6),
            "AI-Powered License Cost Reduction & Security Compliance",
            font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Accent line
add_accent_line(slide, Inches(5), Inches(4.6), Inches(3.333))

# Key stats row
stats = [
    ("34", "Algorithms"),
    ("15-25%", "Cost Savings"),
    ("300-500%", "Annual ROI"),
    ("3-4 Months", "To Production"),
]
for i, (value, label) in enumerate(stats):
    x = Inches(1.5) + Inches(2.7) * i
    add_textbox(slide, x, Inches(5.1), Inches(2.2), Inches(0.6),
                value, font_size=32, color=AZURE_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(5.7), Inches(2.2), Inches(0.4),
                label, font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom bar
bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), Inches(13.333), Pt(4))
bar2.fill.solid()
bar2.fill.fore_color.rgb = AZURE_BLUE
bar2.line.fill.background()

add_textbox(slide, Inches(1), Inches(6.6), Inches(11), Inches(0.4),
            "Built on Azure AI Foundry  |  Microsoft D365 Finance & Operations",
            font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 2: THE PROBLEM
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Pt(6))
bar.fill.solid()
bar.fill.fore_color.rgb = RED_ACCENT
bar.line.fill.background()

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            "The Problem", font_size=36, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
            "D365 FO license management is broken — and it's costing you millions.",
            font_size=18, color=RED_ACCENT, bold=False)

# Problem cards
problems = [
    ("WASTED SPEND", "Over-Licensing", [
        "15-25% of licenses are over-provisioned",
        "Users assigned expensive licenses for read-only work",
        "No visibility into actual usage patterns",
        "Manual reviews happen yearly at best",
    ], RED_ACCENT),
    ("COMPLIANCE RISK", "Security Blind Spots", [
        "SoD violations go undetected for months",
        "Orphaned accounts accumulate silently",
        "No automated audit evidence generation",
        "SOX audits require weeks of manual prep",
    ], AMBER),
    ("OPERATIONAL DRAG", "Manual Processes", [
        "License reviews take 80+ hours per cycle",
        "No data-driven procurement decisions",
        "Reactive instead of proactive management",
        "Role sprawl with 100s of custom roles",
    ], PURPLE),
    ("GROWTH BARRIER", "No Forecasting", [
        "Budget surprises at renewal time",
        "No trend analysis on license needs",
        "Can't model what-if scenarios",
        "Vendor negotiations without usage data",
    ], MID_GRAY),
]

for i, (label, title, lines, color) in enumerate(problems):
    x = Inches(0.5) + Inches(3.15) * i
    add_card(slide, x, Inches(1.8), Inches(2.9), Inches(4.0),
             title, lines, color, label)

# Bottom stat
add_accent_line(slide, Inches(0.8), Inches(6.2), Inches(11.7), RED_ACCENT)
add_textbox(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.8),
            "A 10,000-user organization wastes $3.2M - $5.4M annually on D365 FO license inefficiencies.",
            font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 3: THE SOLUTION
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Pt(6))
bar.fill.solid()
bar.fill.fore_color.rgb = GREEN
bar.line.fill.background()

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            "The Solution", font_size=36, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
            "An AI Agent that continuously optimizes your D365 FO licenses and security posture.",
            font_size=18, color=GREEN)

# Capability cards
caps = [
    ("LICENSE OPTIMIZATION", "Intelligent Cost Reduction", [
        "Detect read-only users for downgrade",
        "Multi-license minority optimization",
        "Cross-role license analysis",
        "Component removal recommendations",
    ], AZURE_BLUE),
    ("SECURITY MONITORING", "Continuous Compliance", [
        "27-rule SoD violation detection",
        "Orphaned account cleanup",
        "Entra ID ↔ D365 FO license sync",
        "Real-time anomaly alerts",
    ], TEAL),
    ("AUTOMATION", "Hands-Off Operations", [
        "New user license wizard",
        "Automated recommendation pipeline",
        "Manager approval workflows",
        "30-day observation mode validation",
    ], PURPLE),
    ("ANALYTICS", "Data-Driven Decisions", [
        "12-month license trend forecasting",
        "Executive cost dashboards",
        "Budget planning & procurement intel",
        "SOX/GDPR/ISO compliance reports",
    ], AMBER),
]

for i, (label, title, lines, color) in enumerate(caps):
    x = Inches(0.5) + Inches(3.15) * i
    add_card(slide, x, Inches(1.8), Inches(2.9), Inches(4.0),
             title, lines, color, label)

add_accent_line(slide, Inches(0.8), Inches(6.2), Inches(11.7), GREEN)
add_textbox(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.8),
            "34 algorithms  |  11 in Phase 1  |  Automated daily/weekly/monthly analysis cycles",
            font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 4: HOW IT WORKS - 5 PHASE LIFECYCLE
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Pt(6))
bar.fill.solid()
bar.fill.fore_color.rgb = AZURE_BLUE
bar.line.fill.background()

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            "How It Works", font_size=36, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
            "Five-Phase AI Agent Lifecycle with Continuous Feedback Loop",
            font_size=18, color=AZURE_LIGHT)

# 5 phase boxes
phases = [
    ("1", "DATA\nACQUISITION", "D365 FO + Azure\nApp Insights\n4 data sources", AZURE_BLUE),
    ("2", "ANALYSIS &\nRECOMMENDATION", "11 algorithms\nConfidence scoring\nPriority ranking", TEAL),
    ("3", "VALIDATION &\nAPPROVAL", "30-day observation\nManager approval\nPeriod-end safeguard", AMBER),
    ("4", "IMPLEMENTATION", "License changes\nUser notifications\nAudit trail", GREEN),
    ("5", "MONITORING &\nIMPROVEMENT", "Post-change tracking\nRollback capability\nCircuit breaker", PURPLE),
]

box_w = Inches(2.2)
box_h = Inches(3.2)
start_x = Inches(0.5)
gap = Inches(0.25)

for i, (num, title, desc, color) in enumerate(phases):
    x = start_x + (box_w + gap) * i

    # Phase box
    card = add_shape(slide, x, Inches(1.7), box_w, box_h, DARK_BG2, color, Pt(2))

    # Phase number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.8), Inches(1.85), Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER

    # Phase title
    add_textbox(slide, x + Inches(0.1), Inches(2.55), box_w - Inches(0.2), Inches(0.8),
                title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Phase desc
    add_textbox(slide, x + Inches(0.15), Inches(3.4), box_w - Inches(0.3), Inches(1.2),
                desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Arrow between boxes (except last)
    if i < 4:
        arrow_x = x + box_w + Pt(2)
        add_textbox(slide, arrow_x, Inches(2.9), Inches(0.3), Inches(0.4),
                    ">>>", font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Feedback loop arrow (text representation)
feedback_shape = add_shape(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.65), DARK_BG2, PURPLE, Pt(1.5))
add_textbox(slide, Inches(0.8), Inches(5.25), Inches(11.7), Inches(0.55),
            "CONTINUOUS FEEDBACK LOOP  <<<  Phase 5 feeds learnings back to Phase 2  |  Algorithm accuracy improves over time  |  Rollback data refines recommendations",
            font_size=13, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

# Web App cross-cutting bar
webapp_shape = add_shape(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.55), DARK_BG2, AZURE_BLUE, Pt(1))
add_textbox(slide, Inches(0.8), Inches(6.15), Inches(11.7), Inches(0.45),
            "WEB APPLICATION (Cross-Cutting)  |  Dashboards  |  Reports  |  Approval Workflows  |  Self-Service Restore  |  Configuration",
            font_size=12, color=AZURE_LIGHT, alignment=PP_ALIGN.CENTER)

# Safeguard gates
add_textbox(slide, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4),
            "8 Safeguard Gates  |  22 Processes  |  RACI Matrix  |  Error Handling & Recovery Flows",
            font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 5: ALGORITHM PORTFOLIO
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Pt(6))
bar.fill.solid()
bar.fill.fore_color.rgb = AZURE_BLUE
bar.line.fill.background()

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            "34 Optimization Algorithms", font_size=36, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
            "Comprehensive coverage across cost, security, behavior, role management, and analytics",
            font_size=18, color=AZURE_LIGHT)

# Category breakdown cards
categories = [
    ("Cost\nOptimization", "12", "7", AZURE_BLUE, [
        "Read-Only User Detection",
        "License Minority Detection",
        "Cross-Role Optimization",
        "Component Removal",
        "Role Splitting",
        "New User License Suggestion",
    ]),
    ("Security &\nCompliance", "9", "2", RED_ACCENT, [
        "SoD Violation Detection",
        "Orphaned Account Cleanup",
        "Entra-D365 License Sync",
        "Privilege Creep Detection",
        "Toxic Combinations",
        "Emergency Account Monitor",
    ]),
    ("User Behavior\nAnalytics", "4", "1", TEAL, [
        "Time-Based Access Analysis",
        "Session Anomaly Detection",
        "Geographic Access Patterns",
        "Usage Pattern Analysis",
    ]),
    ("Role\nManagement", "4", "0", AMBER, [
        "Stale Role Detection",
        "Permission Explosion",
        "Duplicate Role Consolidation",
        "Role Standardization",
    ]),
    ("Advanced\nAnalytics", "5", "1", PURPLE, [
        "License Trend Forecasting",
        "Cost Allocation Engine",
        "What-If Scenario Modeling",
        "ROI Calculator",
        "Utilization Trending",
    ]),
]

card_w = Inches(2.35)
card_h = Inches(4.5)
start_x = Inches(0.4)
gap = Inches(0.18)

for i, (cat_name, total, phase1, color, algos) in enumerate(categories):
    x = start_x + (card_w + gap) * i

    card = add_shape(slide, x, Inches(1.6), card_w, card_h, DARK_BG2, color, Pt(1.5))

    # Category name
    add_textbox(slide, x + Inches(0.15), Inches(1.75), card_w - Inches(0.3), Inches(0.7),
                cat_name, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Stats
    add_textbox(slide, x + Inches(0.15), Inches(2.5), Inches(1.0), Inches(0.5),
                total, font_size=28, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(1.15), Inches(2.5), Inches(1.0), Inches(0.5),
                phase1, font_size=28, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.15), Inches(2.95), Inches(1.0), Inches(0.3),
                "Total", font_size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(1.15), Inches(2.95), Inches(1.0), Inches(0.3),
                "Phase 1", font_size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # Algorithm list
    y = Inches(3.35)
    for algo in algos:
        prefix = "+" if algos.index(algo) < int(phase1) else " "
        c = LIGHT_GRAY if prefix == "+" else MID_GRAY
        add_textbox(slide, x + Inches(0.15), y, card_w - Inches(0.3), Inches(0.25),
                    f"  {algo}", font_size=10, color=c)
        y += Inches(0.22)

# Bottom summary
add_accent_line(slide, Inches(0.8), Inches(6.4), Inches(11.7))
add_textbox(slide, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.5),
            "Phase 1: 11 algorithms (32% of portfolio)  |  Phase 2+: 23 additional algorithms  |  All documented with pseudocode & decision trees",
            font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 6: PHASE 1 TOP 10 ALGORITHMS
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Pt(6))
bar.fill.solid()
bar.fill.fore_color.rgb = GREEN
bar.line.fill.background()

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            "Phase 1: 11 Selected Algorithms", font_size=36, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(8), Inches(0.4),
            "Quick Wins + High ROI + Low-Medium Complexity", font_size=18, color=GREEN)

# Algorithm table
algos_data = [
    ("#1", "Read-Only User Detector", "2.2", "20-40%", "Low", "Highest ROI, immediate impact", AZURE_BLUE),
    ("#2", "License Minority Detection", "2.5", "10-40%", "Med", "Multi-license optimization", AZURE_BLUE),
    ("#3", "Component Removal", "1.4", "5-15%", "Low", "Quick wins, low-hanging fruit", AZURE_BLUE),
    ("#4", "Cross-Role Optimization", "2.6", "10-25%", "Med", "Systemic impact, many users", AZURE_BLUE),
    ("#5", "Role Splitting", "1.3", "10-30%", "Med", "High impact in orgs with overlap", AZURE_BLUE),
    ("#6", "Multi-Role Optimization", "2.4", "5-15%", "Med", "30%+ users have 3+ roles", AZURE_BLUE),
    ("#7", "Orphaned Account Detector", "3.5", "5-10%", "Low", "Security + savings, quick win", RED_ACCENT),
    ("#8", "SoD Violation Detector", "3.1", "Critical", "Med", "SOX compliance mandatory", RED_ACCENT),
    ("#9", "Time-Based Access Analyzer", "5.3", "Security", "Low", "After-hours monitoring", TEAL),
    ("#10", "License Trend Analysis", "4.4", "Strategic", "High", "Budget planning, forecasting", PURPLE),
    ("#11", "New User License Suggestion", "4.7", "5-15%", "Med-Hi", "Prevent over-licensing at source", AZURE_BLUE),
]

# Table header
header_y = Inches(1.55)
cols = [
    (Inches(0.5), Inches(0.5), "#"),
    (Inches(1.0), Inches(2.5), "Algorithm"),
    (Inches(3.5), Inches(0.6), "ID"),
    (Inches(4.1), Inches(1.0), "Savings"),
    (Inches(5.1), Inches(0.7), "Effort"),
    (Inches(5.8), Inches(3.5), "Justification"),
]

# Header row background
hdr_bg = add_shape(slide, Inches(0.4), header_y, Inches(9.0), Inches(0.4), AZURE_BLUE)
for (cx, cw, label) in cols:
    add_textbox(slide, cx, header_y + Pt(2), cw, Inches(0.35),
                label, font_size=12, color=WHITE, bold=True)

# Data rows
for row_i, (rank, name, algo_id, savings, effort, justification, row_color) in enumerate(algos_data):
    y = Inches(1.95) + Inches(0.42) * row_i
    bg_color = DARK_BG2 if row_i % 2 == 0 else DARK_BG
    row_bg = add_shape(slide, Inches(0.4), y, Inches(9.0), Inches(0.4), bg_color)

    add_textbox(slide, Inches(0.5), y + Pt(2), Inches(0.5), Inches(0.35),
                rank, font_size=12, color=row_color, bold=True)
    add_textbox(slide, Inches(1.0), y + Pt(2), Inches(2.5), Inches(0.35),
                name, font_size=12, color=WHITE, bold=False)
    add_textbox(slide, Inches(3.5), y + Pt(2), Inches(0.6), Inches(0.35),
                algo_id, font_size=11, color=MID_GRAY)
    add_textbox(slide, Inches(4.1), y + Pt(2), Inches(1.0), Inches(0.35),
                savings, font_size=12, color=GREEN, bold=True)
    add_textbox(slide, Inches(5.1), y + Pt(2), Inches(0.7), Inches(0.35),
                effort, font_size=11, color=LIGHT_GRAY)
    add_textbox(slide, Inches(5.8), y + Pt(2), Inches(3.5), Inches(0.35),
                justification, font_size=11, color=LIGHT_GRAY)

# Right side - highlight box
highlight_x = Inches(9.8)
highlight_card = add_shape(slide, highlight_x, Inches(1.55), Inches(3.2), Inches(5.5), DARK_BG2, GREEN, Pt(2))

highlight_lines = [
    ("PHASE 1 HIGHLIGHTS", 14, GREEN, True, PP_ALIGN.CENTER),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("7 Cost Optimization", 14, AZURE_BLUE, True, PP_ALIGN.LEFT),
    ("algorithms targeting license waste", 11, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("2 Security & Compliance", 14, RED_ACCENT, True, PP_ALIGN.LEFT),
    ("SoD + orphaned accounts", 11, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("1 Behavior Analytics", 14, TEAL, True, PP_ALIGN.LEFT),
    ("After-hours access detection", 11, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("1 Advanced Analytics", 14, PURPLE, True, PP_ALIGN.LEFT),
    ("12-month license trend forecasting", 11, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("", 10, WHITE, False, PP_ALIGN.LEFT),
    ("25-40% of users optimized", 13, WHITE, True, PP_ALIGN.CENTER),
    ("All data available today", 13, WHITE, True, PP_ALIGN.CENTER),
]
add_multi_text(slide, highlight_x + Inches(0.2), Inches(1.7), Inches(2.8), Inches(5.0), highlight_lines)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 7: BUSINESS IMPACT / ROI
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Pt(6))
bar.fill.solid()
bar.fill.fore_color.rgb = GREEN
bar.line.fill.background()

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            "Business Impact & ROI", font_size=36, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(8), Inches(0.4),
            "Quantified cost savings across organization sizes", font_size=18, color=GREEN)

# Savings table
# Header
table_y = Inches(1.7)
table_cols = [
    (Inches(0.6), Inches(2.5), "Organization Size"),
    (Inches(3.1), Inches(2.0), "Current Annual Spend"),
    (Inches(5.1), Inches(2.2), "Phase 1 Annual Savings"),
    (Inches(7.3), Inches(1.2), "Reduction"),
    (Inches(8.5), Inches(1.5), "Payback Period"),
]

hdr = add_shape(slide, Inches(0.5), table_y, Inches(9.6), Inches(0.45), AZURE_BLUE)
for (cx, cw, label) in table_cols:
    add_textbox(slide, cx, table_y + Pt(3), cw, Inches(0.4),
                label, font_size=13, color=WHITE, bold=True)

rows = [
    ("Small (500 users)", "$1.08M", "$162K - $270K", "15-25%*", "2-3 months"),
    ("Medium (2,000 users)", "$4.32M", "$648K - $1.08M", "15-25%*", "2-3 months"),
    ("Large (10,000 users)", "$21.6M", "$3.2M - $5.4M", "15-25%*", "3-4 months"),
]

for ri, (size, current, savings, pct, payback) in enumerate(rows):
    y = Inches(2.15) + Inches(0.5) * ri
    bg = DARK_BG2 if ri % 2 == 0 else DARK_BG
    add_shape(slide, Inches(0.5), y, Inches(9.6), Inches(0.48), bg)

    add_textbox(slide, Inches(0.6), y + Pt(4), Inches(2.5), Inches(0.4),
                size, font_size=14, color=WHITE, bold=True)
    add_textbox(slide, Inches(3.1), y + Pt(4), Inches(2.0), Inches(0.4),
                current, font_size=14, color=LIGHT_GRAY)
    add_textbox(slide, Inches(5.1), y + Pt(4), Inches(2.2), Inches(0.4),
                savings, font_size=14, color=GREEN, bold=True)
    add_textbox(slide, Inches(7.3), y + Pt(4), Inches(1.2), Inches(0.4),
                pct, font_size=14, color=AZURE_LIGHT, bold=True)
    add_textbox(slide, Inches(8.5), y + Pt(4), Inches(1.5), Inches(0.4),
                payback, font_size=14, color=LIGHT_GRAY)

# Asterisk note
add_textbox(slide, Inches(0.6), Inches(3.85), Inches(9.5), Inches(0.5),
            "*Pending validation of Team Members license form eligibility. Range may increase to 20-35%.",
            font_size=11, color=MID_GRAY)

# ROI highlight cards on right
roi_cards = [
    ("ANNUAL ROI", "300-500%", "Return on investment", GREEN),
    ("PAYBACK", "2-4 Months", "Time to recover costs", AZURE_BLUE),
    ("IMPLEMENTATION", "3-4 Months", "To production ready", TEAL),
]

for i, (label, value, desc, color) in enumerate(roi_cards):
    x = Inches(10.5)
    y = Inches(1.7) + Inches(1.05) * i
    card = add_shape(slide, x, y, Inches(2.5), Inches(0.9), DARK_BG2, color, Pt(2))
    add_textbox(slide, x + Inches(0.15), y + Pt(5), Inches(2.2), Inches(0.3),
                label, font_size=10, color=color, bold=True)
    add_textbox(slide, x + Inches(0.15), y + Pt(22), Inches(2.2), Inches(0.3),
                value, font_size=22, color=WHITE, bold=True)
    add_textbox(slide, x + Inches(0.15), y + Pt(50), Inches(2.2), Inches(0.25),
                desc, font_size=10, color=MID_GRAY)

# Operational benefits
add_accent_line(slide, Inches(0.6), Inches(4.4), Inches(12.2))
add_textbox(slide, Inches(0.6), Inches(4.6), Inches(12), Inches(0.5),
            "Operational Benefits", font_size=22, color=WHITE, bold=True)

ops_benefits = [
    ("80% Reduction", "in manual review time", AZURE_BLUE),
    ("90% Reduction", "in audit preparation time", GREEN),
    ("25-40% of Users", "optimized with license changes", TEAL),
    ("SOX Foundation", "compliance established from day 1", PURPLE),
]

for i, (stat, desc, color) in enumerate(ops_benefits):
    x = Inches(0.6) + Inches(3.15) * i
    card = add_shape(slide, x, Inches(5.15), Inches(2.9), Inches(0.9), DARK_BG2, color, Pt(1))
    add_textbox(slide, x + Inches(0.2), Inches(5.2), Inches(2.5), Inches(0.4),
                stat, font_size=18, color=color, bold=True)
    add_textbox(slide, x + Inches(0.2), Inches(5.6), Inches(2.5), Inches(0.3),
                desc, font_size=12, color=LIGHT_GRAY)

# Bottom KPIs
add_textbox(slide, Inches(0.6), Inches(6.3), Inches(12.2), Inches(0.8),
            "KPIs: 15-25% cost reduction (+3-8% with Entra sync)  |  25-40% users optimized  |  < 24h data freshness  |  > 80% user adoption  |  > 70% recommendation acceptance  |  100% audit readiness",
            font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 8: ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Pt(6))
bar.fill.solid()
bar.fill.fore_color.rgb = AZURE_BLUE
bar.line.fill.background()

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            "Architecture", font_size=36, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(8), Inches(0.4),
            "Built on Azure AI Foundry + Azure AI Agent Service", font_size=18, color=AZURE_LIGHT)

# Data Sources (left column)
add_textbox(slide, Inches(0.5), Inches(1.55), Inches(3.0), Inches(0.4),
            "DATA SOURCES", font_size=14, color=AZURE_BLUE, bold=True)

sources = [
    ("D365 FO Security Config", "700K+ records via OData"),
    ("D365 FO User-Role Data", "200K records via OData"),
    ("Azure App Insights", "10M events/day streaming"),
    ("D365 FO Audit Logs", "Real-time change tracking"),
    ("Microsoft Graph API", "Entra ID licenses (optional)"),
]
for i, (name, desc) in enumerate(sources):
    y = Inches(1.95) + Inches(0.55) * i
    card = add_shape(slide, Inches(0.5), y, Inches(3.0), Inches(0.5), DARK_BG2, AZURE_BLUE, Pt(1))
    add_textbox(slide, Inches(0.65), y + Pt(2), Inches(2.7), Inches(0.25),
                name, font_size=12, color=WHITE, bold=True)
    add_textbox(slide, Inches(0.65), y + Pt(18), Inches(2.7), Inches(0.2),
                desc, font_size=10, color=MID_GRAY)

# Arrow
add_textbox(slide, Inches(3.6), Inches(2.7), Inches(0.5), Inches(0.5),
            ">>>", font_size=20, color=AZURE_LIGHT, bold=True)

# Azure AI Agent (center - 6 layers)
add_textbox(slide, Inches(4.2), Inches(1.55), Inches(5.0), Inches(0.4),
            "AZURE AI AGENT (6 Layers)", font_size=14, color=TEAL, bold=True)

layers = [
    ("1. Data Ingestion", "OData connectors, Event Hubs, validation", AZURE_BLUE),
    ("2. Algorithm Engine", "34 algorithms, parallel processing", TEAL),
    ("3. Recommendation Engine", "Confidence scoring, priority ranking", GREEN),
    ("4. Scheduling & Automation", "Daily/weekly/monthly + event-triggered", AMBER),
    ("5. API & Integration", "REST APIs, WebSocket, batch export", PURPLE),
    ("6. Audit & Compliance", "Immutable logs, data lineage, 7yr retention", RED_ACCENT),
]

for i, (name, desc, color) in enumerate(layers):
    y = Inches(1.95) + Inches(0.55) * i
    card = add_shape(slide, Inches(4.2), y, Inches(5.0), Inches(0.5), DARK_BG2, color, Pt(1))
    add_textbox(slide, Inches(4.35), y + Pt(2), Inches(2.0), Inches(0.25),
                name, font_size=12, color=color, bold=True)
    add_textbox(slide, Inches(6.4), y + Pt(2), Inches(2.7), Inches(0.25),
                desc, font_size=10, color=LIGHT_GRAY)

# Arrow
add_textbox(slide, Inches(9.3), Inches(2.7), Inches(0.5), Inches(0.5),
            ">>>", font_size=20, color=AZURE_LIGHT, bold=True)

# Web Application (right column)
add_textbox(slide, Inches(9.9), Inches(1.55), Inches(3.0), Inches(0.4),
            "WEB APPLICATION", font_size=14, color=GREEN, bold=True)

web_items = [
    ("Executive Dashboard", "Cost savings, trends, health"),
    ("Recommendation Center", "Review, approve, track changes"),
    ("Compliance Reports", "SOX, GDPR, ISO evidence"),
    ("Self-Service Portal", "User restore, status tracking"),
]
for i, (name, desc) in enumerate(web_items):
    y = Inches(1.95) + Inches(0.55) * i
    card = add_shape(slide, Inches(9.9), y, Inches(3.0), Inches(0.5), DARK_BG2, GREEN, Pt(1))
    add_textbox(slide, Inches(10.05), y + Pt(2), Inches(2.7), Inches(0.25),
                name, font_size=12, color=WHITE, bold=True)
    add_textbox(slide, Inches(10.05), y + Pt(18), Inches(2.7), Inches(0.2),
                desc, font_size=10, color=MID_GRAY)

# Azure services footer
add_accent_line(slide, Inches(0.5), Inches(5.3), Inches(12.3))

azure_services = [
    ("Compute", "Azure Functions\nAzure Batch", AZURE_BLUE),
    ("Storage", "Azure SQL\nData Lake Gen2\nCosmos DB", TEAL),
    ("Integration", "Event Hubs\nAPI Management\nLogic Apps", GREEN),
    ("Monitoring", "App Insights\nAzure Monitor\nLog Analytics", AMBER),
    ("Security", "Key Vault\nAzure AD\nAzure Policy", RED_ACCENT),
]

for i, (name, services, color) in enumerate(azure_services):
    x = Inches(0.5) + Inches(2.55) * i
    card = add_shape(slide, x, Inches(5.5), Inches(2.35), Inches(1.5), DARK_BG2, color, Pt(1))
    add_textbox(slide, x + Inches(0.15), Inches(5.55), Inches(2.0), Inches(0.3),
                name, font_size=13, color=color, bold=True)
    add_textbox(slide, x + Inches(0.15), Inches(5.85), Inches(2.0), Inches(1.0),
                services, font_size=10, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 9: SECURITY & COMPLIANCE
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Pt(6))
bar.fill.solid()
bar.fill.fore_color.rgb = RED_ACCENT
bar.line.fill.background()

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            "Security & Compliance", font_size=36, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
            "Enterprise-grade security with built-in compliance frameworks",
            font_size=18, color=RED_ACCENT)

# SoD section
sod_card = add_shape(slide, Inches(0.5), Inches(1.6), Inches(6.0), Inches(3.0), DARK_BG2, RED_ACCENT, Pt(2))
add_textbox(slide, Inches(0.7), Inches(1.7), Inches(5.5), Inches(0.4),
            "Separation of Duties (SoD) Detection", font_size=18, color=RED_ACCENT, bold=True)

sod_categories = [
    "Accounts Payable (7 rules) - Invoice creation vs. payment approval",
    "Accounts Receivable (4 rules) - Customer credit vs. collections",
    "General Ledger (4 rules) - Journal entry vs. posting",
    "Procurement (4 rules) - Purchase orders vs. vendor management",
    "Fixed Assets (3 rules) - Acquisition vs. disposal",
    "Inventory (3 rules) - Adjustment vs. counting",
    "System Admin (2 rules) - Security config vs. business roles",
]
for i, cat in enumerate(sod_categories):
    add_textbox(slide, Inches(0.8), Inches(2.2) + Inches(0.3) * i, Inches(5.5), Inches(0.3),
                cat, font_size=11, color=LIGHT_GRAY)

add_textbox(slide, Inches(0.8), Inches(4.35), Inches(5.5), Inches(0.3),
            "27 rules across 7 categories  |  Industry-standard conflict matrix",
            font_size=12, color=WHITE, bold=True)

# Compliance frameworks (right)
frameworks = [
    ("SOX Section 404", "Access control documentation\nRole assignment history\nChange management audit trail\nAutomated evidence generation", RED_ACCENT),
    ("GDPR", "PII data protection & encryption\nData access documentation\nRight to erasure compliance\nData breach notification", AZURE_BLUE),
    ("ISO 27001", "Security monitoring evidence\nAccess review audit logs\nIncident response documentation\nRisk assessment results", TEAL),
]

for i, (name, details, color) in enumerate(frameworks):
    y = Inches(1.6) + Inches(1.05) * i
    card = add_shape(slide, Inches(6.8), y, Inches(6.0), Inches(0.95), DARK_BG2, color, Pt(1.5))
    add_textbox(slide, Inches(7.0), y + Pt(5), Inches(1.8), Inches(0.3),
                name, font_size=14, color=color, bold=True)
    add_textbox(slide, Inches(8.8), y + Pt(5), Inches(3.8), Inches(0.85),
                details, font_size=10, color=LIGHT_GRAY)

# Safeguard gates
add_accent_line(slide, Inches(0.5), Inches(4.8), Inches(12.3), RED_ACCENT)
add_textbox(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.4),
            "8 Safeguard Gates Protecting Every Phase", font_size=20, color=WHITE, bold=True)

gates = [
    ("G1", "Data validation\nthresholds", AZURE_BLUE),
    ("G2", "Minimum confidence\nscores", TEAL),
    ("G3", "30-day observation\nmode pass", AMBER),
    ("G4", "Manager approval\nrequired", GREEN),
    ("G5", "Period-end\nfreeze check", PURPLE),
    ("G6", "Circuit breaker\nmonitoring", RED_ACCENT),
    ("G7", "Rollback SLA\nenforcement", AZURE_BLUE),
    ("G8", "Escalation tier\nvalidation", TEAL),
]

for i, (label, desc, color) in enumerate(gates):
    x = Inches(0.5) + Inches(1.58) * i
    card = add_shape(slide, x, Inches(5.5), Inches(1.45), Inches(1.2), DARK_BG2, color, Pt(1))
    # Gate label
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.45), Inches(5.55), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER

    add_textbox(slide, x + Inches(0.1), Inches(6.1), Inches(1.25), Inches(0.55),
                desc, font_size=9, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 10: IMPLEMENTATION ROADMAP
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Pt(6))
bar.fill.solid()
bar.fill.fore_color.rgb = AZURE_BLUE
bar.line.fill.background()

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            "Implementation Roadmap", font_size=36, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(8), Inches(0.4),
            "14 weeks to production across 4 sprints", font_size=18, color=AZURE_LIGHT)

# Timeline sprints
sprints = [
    ("SPRINT 1", "Foundation & Quick Wins", "Weeks 1-3", [
        "Read-Only User Detector (2.2)",
        "Component Removal (1.4)",
        "Orphaned Account Detector (3.5)",
        "Time-Based Access Analyzer (5.3)",
    ], "$150K-$750K savings", AZURE_BLUE),
    ("SPRINT 2", "Advanced Optimization", "Weeks 4-7", [
        "License Minority Detection (2.5)",
        "Cross-Role Optimization (2.6)",
        "Role Splitting Recommender (1.3)",
        "",
    ], "$230K-$1M savings", TEAL),
    ("SPRINT 3", "Compliance & Analytics", "Weeks 8-12", [
        "Multi-Role Optimization (2.4)",
        "SoD Violation Detector (3.1)",
        "License Trend Analysis (4.4)",
        "",
    ], "$280K-$1.2M savings", AMBER),
    ("SPRINT 4", "Integration & Launch", "Weeks 12-14", [
        "End-to-end integration",
        "User acceptance testing (UAT)",
        "Performance optimization",
        "Production deployment",
    ], "PRODUCTION READY", GREEN),
]

card_w = Inches(3.0)
card_h = Inches(4.2)
start_x = Inches(0.35)
gap = Inches(0.15)

for i, (sprint_label, title, timeline, algos, impact, color) in enumerate(sprints):
    x = start_x + (card_w + gap) * i

    card = add_shape(slide, x, Inches(1.6), card_w, card_h, DARK_BG2, color, Pt(2))

    # Sprint label
    label_bg = add_shape(slide, x, Inches(1.6), card_w, Inches(0.35), color)
    add_textbox(slide, x + Inches(0.15), Inches(1.63), card_w - Inches(0.3), Inches(0.3),
                f"{sprint_label}  |  {timeline}", font_size=12, color=WHITE, bold=True)

    # Title
    add_textbox(slide, x + Inches(0.15), Inches(2.05), card_w - Inches(0.3), Inches(0.4),
                title, font_size=16, color=WHITE, bold=True)

    # Algorithms
    for j, algo in enumerate(algos):
        if algo:
            add_textbox(slide, x + Inches(0.2), Inches(2.55) + Inches(0.3) * j,
                        card_w - Inches(0.4), Inches(0.3),
                        f"  {algo}", font_size=11, color=LIGHT_GRAY)

    # Impact badge
    impact_y = Inches(3.85) + Inches(0.8)
    impact_bg = add_shape(slide, x + Inches(0.15), impact_y, card_w - Inches(0.3), Inches(0.4), color)
    add_textbox(slide, x + Inches(0.2), impact_y + Pt(2), card_w - Inches(0.4), Inches(0.35),
                impact, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Arrow between sprints
    if i < 3:
        arrow_x = x + card_w + Pt(2)
        add_textbox(slide, arrow_x, Inches(3.3), Inches(0.2), Inches(0.4),
                    ">", font_size=18, color=MID_GRAY)

# Bottom timeline bar
timeline_y = Inches(6.1)
add_shape(slide, Inches(0.5), timeline_y, Inches(12.3), Inches(0.5), DARK_BG2, AZURE_BLUE, Pt(1))

# Week markers
weeks = [
    ("Week 1", Inches(0.7)),
    ("Week 3", Inches(3.3)),
    ("Week 7", Inches(6.0)),
    ("Week 12", Inches(9.0)),
    ("Week 14", Inches(11.5)),
]
for label, wx in weeks:
    add_textbox(slide, wx, timeline_y + Pt(5), Inches(1.2), Inches(0.35),
                label, font_size=11, color=AZURE_LIGHT, bold=True)

# Progress bar fill
add_shape(slide, Inches(0.5), timeline_y, Inches(12.3), Pt(4), GREEN)

add_textbox(slide, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
            "3-4 developers  |  Agile sprints  |  CI/CD pipeline  |  Blue-Green deployment  |  Infrastructure as Code (Terraform/Bicep)",
            font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 11: WHY THIS AGENT / DIFFERENTIATORS
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Pt(6))
bar.fill.solid()
bar.fill.fore_color.rgb = PURPLE
bar.line.fill.background()

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            "Why This Agent", font_size=36, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
            "Purpose-built for D365 FO with deep domain expertise",
            font_size=18, color=PURPLE)

# Differentiator cards (2x3 grid)
diffs = [
    ("D365 FO NATIVE", "Purpose-Built Intelligence", [
        "Understands D365 FO licensing model natively",
        "Knows that 'highest license wins' across roles",
        "Leverages read-only access rules for downgrades",
        "Handles Team Members / Operations / Commerce tiers",
    ], AZURE_BLUE),
    ("AI-POWERED", "Not Just Rules — Intelligence", [
        "Confidence scoring on every recommendation",
        "Pattern recognition across user populations",
        "Continuous learning from rollback feedback",
        "Trend analysis and 12-month forecasting",
    ], TEAL),
    ("SAFETY-FIRST", "8 Gates, Zero Disruption", [
        "30-day observation mode before any changes",
        "Manager approval workflow required",
        "Circuit breaker auto-disables on high rollback",
        "Period-end safeguard freezes during close",
    ], GREEN),
    ("COMPLIANCE-READY", "Audit From Day 1", [
        "SOX Section 404 evidence auto-generated",
        "27-rule SoD conflict detection matrix",
        "Immutable audit trail with 7-year retention",
        "GDPR and ISO 27001 reporting built-in",
    ], RED_ACCENT),
    ("SCALABLE", "500 to 50,000 Users", [
        "Horizontal scaling via Azure infrastructure",
        "Delta sync reduces OData load by 95%",
        "Processes 10M telemetry events per day",
        "API handles 1,000+ requests per second",
    ], AMBER),
    ("PHASED DELIVERY", "Value From Week 1", [
        "Phase 1: 11 algorithms, immediate savings",
        "Phase 2: 23 more algorithms ready to go",
        "Modular architecture — add algorithms easily",
        "ROI starts in first sprint, not after go-live",
    ], PURPLE),
]

for i, (label, title, lines, color) in enumerate(diffs):
    col = i % 3
    row = i // 3
    x = Inches(0.4) + Inches(4.2) * col
    y = Inches(1.6) + Inches(2.55) * row
    add_card(slide, x, y, Inches(3.95), Inches(2.35), title, lines, color, label)

# Bottom tagline
add_accent_line(slide, Inches(0.8), Inches(6.8), Inches(11.7), PURPLE)
add_textbox(slide, Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.5),
            "34 algorithms  |  23 documents of requirements  |  3,000+ pages of specifications  |  Built to enterprise standards",
            font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 12: NEXT STEPS / CTA
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Pt(6))
bar.fill.solid()
bar.fill.fore_color.rgb = AZURE_BLUE
bar.line.fill.background()

add_textbox(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
            "Ready to Optimize?", font_size=48, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(2.2), Inches(9), Inches(0.6),
            "Start saving 15-25% on D365 FO license costs in 3-4 months.",
            font_size=22, color=AZURE_LIGHT, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5), Inches(3.0), Inches(3.333))

# Next steps cards
steps = [
    ("1", "Stakeholder Review", "Present Phase 1 selection\nValidate priorities\nConfirm budget", AZURE_BLUE),
    ("2", "Technical Feasibility", "Validate data access\nArchitecture planning\nAzure service selection", TEAL),
    ("3", "Sprint Planning", "Team allocation\nDetailed sprint breakdown\nKickoff meeting", GREEN),
    ("4", "Development Begins", "Sprint 1 launches\nFirst results in 2-3 weeks\nROI from day 1", AMBER),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.5) + Inches(3.15) * i
    card = add_shape(slide, x, Inches(3.5), Inches(2.9), Inches(2.0), DARK_BG2, color, Pt(2))

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.1), Inches(3.6), Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER

    add_textbox(slide, x + Inches(0.2), Inches(4.25), Inches(2.5), Inches(0.35),
                title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), Inches(4.65), Inches(2.5), Inches(0.8),
                desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom bar with key stats
add_accent_line(slide, Inches(0.8), Inches(5.8), Inches(11.7))

summary_stats = [
    ("34 Algorithms", "11 in Phase 1"),
    ("15-25% Savings", "300-500% ROI"),
    ("3-4 Month", "Implementation"),
    ("2-4 Month", "Payback Period"),
]

for i, (big, small) in enumerate(summary_stats):
    x = Inches(1.2) + Inches(3.0) * i
    add_textbox(slide, x, Inches(5.95), Inches(2.5), Inches(0.4),
                big, font_size=20, color=AZURE_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(6.35), Inches(2.5), Inches(0.3),
                small, font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Final line
bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), Inches(13.333), Pt(4))
bar2.fill.solid()
bar2.fill.fore_color.rgb = AZURE_BLUE
bar2.line.fill.background()

add_textbox(slide, Inches(1), Inches(6.7), Inches(11), Inches(0.4),
            "D365 FO License & Security Optimization Agent  |  Built on Azure AI Foundry",
            font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_path = "/home/user/projects/work/D365FOLicenseAgent-v1/D365-FO-License-Agent-Sales-Deck.pptx"
prs.save(output_path)
print(f"Sales deck saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
